import os
from flask import Flask, request, render_template, send_file
from werkzeug.utils import secure_filename
import logging
from pptx import Presentation
from flask import send_file
from azure.storage.blob import BlobServiceClient, BlobClient, ContainerClient
from azure.core.exceptions import ResourceNotFoundError

# ... (Rest Ihres Codes)

connect_str = "DefaultEndpointsProtocol=https;AccountName=speicherstorage;AccountKey=C17G8ekAUuM0D32V3kDC2bRWEiLP1872EvOGDWx//AiJOEj/bmIM7PTNItb7nbnhZodqOu3LISzY+ASt8Kpfjg==;EndpointSuffix=core.windows.net"
blob_service_client = BlobServiceClient.from_connection_string(connect_str)

logging.basicConfig(level=logging.INFO)
app = Flask(__name__)
app.secret_key = 'ein_zufälliger_schlüssel'
UPLOAD_FOLDER = r'C:\Users\admin\Desktop\App_OCR\app _pii\neue'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

BRAND_NAME_MAPPINGS = {
    'Makrosoft': 'Microsoft',
    'Lenevo': 'Lenovo',
    'Sorny': 'Sony',
    'Adibas': 'Adidas',
    'Sansung': 'Samsung'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'ppt', 'pptx'}

def analyze_pptx_and_find_corrections(pptx_path):
    prs = Presentation(pptx_path)
    corrections = {}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                for fake_name, real_name in BRAND_NAME_MAPPINGS.items():
                    if fake_name in text:
                        corrections[fake_name] = real_name
    return corrections

def save_results_to_txt(corrections, filepath):
    with open(filepath, 'w') as file:
        for fake, real in corrections.items():
            file.write(f"{fake}: {real}\n")

#################
# Upload- und Download-Funktionen für Azure Blob Storage
def upload_file_to_blob(file_name, blob_name, container_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    with open(file_name, "rb") as data:
        blob_client.upload_blob(data)

def download_file_from_blob(blob_name, download_file_path, container_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    with open(download_file_path, "wb") as download_file:
        download_file.write(blob_client.download_blob().readall())
        ###########
        
        
def blob_exists(container_name, blob_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    try:
        blob_client.get_blob_properties()
        return True
    except ResourceNotFoundError:
        return False

############

def generate_unique_blob_name(container_name, original_name):
    count = 1
    unique_name = f"{original_name}_{count}"
    while blob_exists(container_name, unique_name):
        count += 1
        unique_name = f"{original_name}_{count}"
    return unique_name

##################
@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            saved_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(saved_path)

            if filename.endswith('.pptx'):
                corrections = analyze_pptx_and_find_corrections(saved_path)

                container_name = "hasan"  # Ihr tatsächlicher Container-Name

                unique_blob_name = generate_unique_blob_name(container_name, filename)
                txt_filename = f'Ergebnisse_{unique_blob_name}.txt'
                temp_txt_path = os.path.join(app.config['UPLOAD_FOLDER'], txt_filename)

                # Speichern der Ergebnisse in der Textdatei
                save_results_to_txt(corrections, temp_txt_path)

                # Hochladen der Dateien in Azure Blob Storage
                upload_file_to_blob(temp_txt_path, txt_filename, container_name)
                upload_file_to_blob(saved_path, unique_blob_name, container_name)

                # Aufräumen: Originaldatei entfernen
                os.remove(saved_path)

                return render_template('results.html', message=corrections, filename=txt_filename)
            else:
                os.remove(saved_path)
                return 'Ungültiges Dateiformat', 400

    return render_template('upload.html')



@app.route('/download_results', methods=['POST'])
def download_results():
    filename = request.form['filename']
    if filename:
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        try:
            return send_file(filepath, as_attachment=True, attachment_filename=filename)
        except Exception as e:
            logging.error(f"An error occurred: {e}")
            return str(e), 500
    else:
        return 'Fehler: Kein Dateiname angegeben', 400


if __name__ == "__main__":
    app.run(debug=True)
